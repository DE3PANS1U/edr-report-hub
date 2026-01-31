"""
Detailed chart data analysis - check if charts have actual data values
"""
from pptx import Presentation

prs = Presentation("EDR-Weekly-Incident-Report_2026-01-31.pptx")

print("=" * 80)
print("DETAILED CHART DATA ANALYSIS")
print("=" * 80)

for slide_idx in [1, 2, 3, 4, 5]:  # Slides 2-6 (0-indexed: 1-5)
    slide = prs.slides[slide_idx]
    print(f"\n{'=' * 80}")
    print(f"SLIDE {slide_idx + 1}")
    print(f"{'=' * 80}")
    
    chart_num = 0
    for shape in slide.shapes:
        if hasattr(shape, "has_chart") and shape.has_chart:
            chart_num += 1
            chart = shape.chart
            
            print(f"\nChart {chart_num}: {chart.chart_title.text_frame.text if chart.has_title else 'No title'}")
            print(f"  Type: {chart.chart_type}")
            
            # Try to get chart data
            try:
                print(f"  Series count: {len(chart.series)}")
                
                for series_idx, series in enumerate(chart.series):
                    print(f"  Series {series_idx + 1}: {series.name}")
                    
                    # Get values
                    try:
                        values = list(series.values)
                        print(f"    Values: {values}")
                        print(f"    Sum: {sum(values)}")
                        print(f"    Non-zero count: {sum(1 for v in values if v > 0)}")
                    except Exception as e:
                        print(f"    ERROR getting values: {e}")
                
                # Get categories
                try:
                    if hasattr(chart, 'plots') and len(chart.plots) > 0:
                        plot = chart.plots[0]
                        if hasattr(plot, 'categories'):
                            categories = [cat for cat in plot.categories]
                            print(f"  Categories: {categories}")
                except Exception as e:
                    print(f"  ERROR getting categories: {e}")
                    
            except Exception as e:
                print(f"  ERROR analyzing chart: {e}")

print("\n" + "=" * 80)
print("ANALYSIS COMPLETE")
print("=" * 80)
