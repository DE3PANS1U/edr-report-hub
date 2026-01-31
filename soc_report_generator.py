"""
SOC Weekly Report Generator
Generates professional PowerPoint reports from EDR alert data
"""

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.dml.color import RGBColor
from datetime import datetime
import os


# Corporate SOC Color Theme
COLORS = {
    'critical': RGBColor(220, 53, 69),      # Red
    'high': RGBColor(255, 133, 27),         # Orange
    'medium': RGBColor(255, 193, 7),        # Amber
    'low': RGBColor(40, 167, 69),           # Green
    'closed': RGBColor(40, 167, 69),        # Green
    'wip': RGBColor(255, 193, 7),           # Amber
    'true_positive': RGBColor(220, 53, 69), # Red
    'false_positive': RGBColor(40, 167, 69),# Green
    'background': RGBColor(248, 249, 250),  # Light Gray
    'text': RGBColor(33, 37, 41),           # Dark Gray
    'accent': RGBColor(0, 123, 255)         # Blue
}

SEVERITY_ORDER = ['Critical', 'High', 'Medium', 'Low']
STATUS_ORDER = ['closed', 'WIP']
EFFICIENCY_ORDER = ['True Positive', 'False Positive']


def read_excel_data(file_path):
    """Read EDR alert data from Excel file"""
    try:
        df = pd.read_excel(file_path)
        
        # Ensure required columns exist
        required_columns = ['Entity', 'Alert Severity', 'Alert Status', 
                          'Alert Efficiency', 'Alert Trend']
        
        missing_columns = [col for col in required_columns if col not in df.columns]
        if missing_columns:
            raise ValueError(f"Missing required columns: {missing_columns}")
        
        # Convert Alert Trend to datetime
        df['Alert Trend'] = pd.to_datetime(df['Alert Trend'])
        
        return df
    except Exception as e:
        raise Exception(f"Error reading Excel file: {str(e)}")


def add_title_slide(prs, title, subtitle):
    """Add title slide to presentation"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    # Format title
    title_frame = title_shape.text_frame
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLORS['accent']
    
    # Format subtitle
    subtitle_frame = subtitle_shape.text_frame
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    subtitle_frame.paragraphs[0].font.color.rgb = COLORS['text']


def create_clustered_bar_chart(slide, chart_data, title, left, top, width, height, categories_colors):
    """Create a clustered bar chart"""
    chart_placeholder = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data
    )
    
    chart = chart_placeholder.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(18)
    chart.chart_title.text_frame.paragraphs[0].font.bold = True
    
    # Format legend
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = Pt(11)
    
    # Apply colors to series
    for idx, series in enumerate(chart.series):
        if idx < len(categories_colors):
            category_name = list(categories_colors.keys())[idx]
            color = categories_colors[category_name]
            fill = series.format.fill
            fill.solid()
            fill.fore_color.rgb = color
    
    return chart


def create_line_chart(slide, chart_data, title, left, top, width, height, series_colors):
    """Create a line chart with markers"""
    chart_placeholder = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, left, top, width, height, chart_data
    )
    
    chart = chart_placeholder.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(18)
    chart.chart_title.text_frame.paragraphs[0].font.bold = True
    
    # Format legend
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = Pt(10)
    
    # Apply colors to series
    for idx, series in enumerate(chart.series):
        if idx < len(series_colors):
            color = series_colors[idx]
            # Line color
            line = series.format.line
            line.color.rgb = color
            line.width = Pt(2.5)
            
            # Marker
            series.marker.style = 2  # Circle marker
            series.marker.size = 7
    
    return chart


def add_table_to_slide(slide, data, left, top, width, height):
    """Add a formatted table to slide"""
    rows, cols = len(data), len(data[0])
    
    table_placeholder = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_placeholder.table
    
    # Populate table
    for i, row_data in enumerate(data):
        for j, cell_value in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = str(cell_value)
            
            # Format cell
            cell.text_frame.paragraphs[0].font.size = Pt(11)
            
            # Header row formatting
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['accent']
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            else:
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return table


def generate_severity_slide(prs, df):
    """Slide 1: Alert Severity Distribution"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    
    # Add slide title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "Alert Severity Distribution"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLORS['accent']
    
    # Prepare data - Severity distribution per entity
    severity_df = df.groupby(['Entity', 'Alert Severity']).size().unstack(fill_value=0)
    
    # Reorder columns to match severity order
    severity_df = severity_df.reindex(columns=SEVERITY_ORDER, fill_value=0)
    
    # Create chart data
    chart_data = CategoryChartData()
    chart_data.categories = severity_df.index.tolist()
    
    for severity in SEVERITY_ORDER:
        if severity in severity_df.columns:
            chart_data.add_series(severity, severity_df[severity].tolist())
    
    # Colors for severity
    severity_colors = {
        'Critical': COLORS['critical'],
        'High': COLORS['high'],
        'Medium': COLORS['medium'],
        'Low': COLORS['low']
    }
    
    # Add chart
    create_clustered_bar_chart(
        slide, chart_data, 
        "Severity Distribution by Entity",
        Inches(0.5), Inches(1.2), Inches(9), Inches(3.5),
        severity_colors
    )
    
    # Create summary table
    table_data = [['Entity', 'Critical', 'High', 'Medium', 'Low', 'Total']]
    
    for entity in severity_df.index:
        row = [entity]
        for severity in SEVERITY_ORDER:
            row.append(int(severity_df.loc[entity, severity]))
        row.append(int(severity_df.loc[entity].sum()))
        table_data.append(row)
    
    # Add totals row
    totals_row = ['Total']
    for severity in SEVERITY_ORDER:
        totals_row.append(int(severity_df[severity].sum()))
    totals_row.append(int(df.shape[0]))
    table_data.append(totals_row)
    
    # Add table
    add_table_to_slide(slide, table_data, Inches(0.5), Inches(5.0), Inches(9), Inches(2))


def generate_status_slide(prs, df):
    """Slide 2: Alert Status Overview"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    # Add slide title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "Alert Status Overview"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLORS['accent']
    
    # Prepare data - Status distribution per entity
    status_df = df.groupby(['Entity', 'Alert Status']).size().unstack(fill_value=0)
    
    # Ensure both statuses exist
    for status in STATUS_ORDER:
        if status not in status_df.columns:
            status_df[status] = 0
    
    status_df = status_df[STATUS_ORDER]
    
    # Create chart data
    chart_data = CategoryChartData()
    chart_data.categories = status_df.index.tolist()
    
    chart_data.add_series('Closed', status_df['closed'].tolist())
    chart_data.add_series('WIP', status_df['WIP'].tolist())
    
    # Colors for status
    status_colors = {
        'closed': COLORS['closed'],
        'WIP': COLORS['wip']
    }
    
    # Add chart
    create_clustered_bar_chart(
        slide, chart_data,
        "Status Distribution by Entity",
        Inches(0.5), Inches(1.2), Inches(9), Inches(3.5),
        status_colors
    )
    
    # Create summary table
    table_data = [['Entity', 'Closed', 'WIP', 'Total']]
    
    for entity in status_df.index:
        row = [entity, int(status_df.loc[entity, 'closed']), 
               int(status_df.loc[entity, 'WIP']), 
               int(status_df.loc[entity].sum())]
        table_data.append(row)
    
    # Add totals row
    totals_row = ['Total', int(status_df['closed'].sum()), 
                  int(status_df['WIP'].sum()), int(df.shape[0])]
    table_data.append(totals_row)
    
    # Add table
    add_table_to_slide(slide, table_data, Inches(2.5), Inches(5.0), Inches(5), Inches(2))


def generate_efficiency_slide(prs, df):
    """Slide 3: Alert Efficiency Analysis"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    # Add slide title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "Alert Efficiency Analysis"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLORS['accent']
    
    # Prepare data - Efficiency distribution per entity
    efficiency_df = df.groupby(['Entity', 'Alert Efficiency']).size().unstack(fill_value=0)
    
    # Ensure both efficiency types exist
    for eff in EFFICIENCY_ORDER:
        if eff not in efficiency_df.columns:
            efficiency_df[eff] = 0
    
    efficiency_df = efficiency_df[EFFICIENCY_ORDER]
    
    # Create chart data
    chart_data = CategoryChartData()
    chart_data.categories = efficiency_df.index.tolist()
    
    chart_data.add_series('True Positive', efficiency_df['True Positive'].tolist())
    chart_data.add_series('False Positive', efficiency_df['False Positive'].tolist())
    
    # Colors for efficiency
    efficiency_colors = {
        'True Positive': COLORS['true_positive'],
        'False Positive': COLORS['false_positive']
    }
    
    # Add chart
    create_clustered_bar_chart(
        slide, chart_data,
        "Efficiency Analysis by Entity",
        Inches(0.5), Inches(1.2), Inches(9), Inches(3.5),
        efficiency_colors
    )
    
    # Create summary table
    table_data = [['Entity', 'True Positive', 'False Positive', 'Total', 'Accuracy %']]
    
    for entity in efficiency_df.index:
        tp = int(efficiency_df.loc[entity, 'True Positive'])
        fp = int(efficiency_df.loc[entity, 'False Positive'])
        total = tp + fp
        accuracy = (tp / total * 100) if total > 0 else 0
        row = [entity, tp, fp, total, f"{accuracy:.1f}%"]
        table_data.append(row)
    
    # Add totals row
    total_tp = int(efficiency_df['True Positive'].sum())
    total_fp = int(efficiency_df['False Positive'].sum())
    total_all = total_tp + total_fp
    total_accuracy = (total_tp / total_all * 100) if total_all > 0 else 0
    totals_row = ['Total', total_tp, total_fp, total_all, f"{total_accuracy:.1f}%"]
    table_data.append(totals_row)
    
    # Add table
    add_table_to_slide(slide, table_data, Inches(1.5), Inches(5.0), Inches(7), Inches(2))


def generate_trend_slide(prs, df):
    """Slide 4: Alert Trend Analysis"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    # Add slide title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "Alert Trend Analysis (Editable)"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLORS['accent']
    
    # Prepare data - Daily trend per entity
    df['Date'] = df['Alert Trend'].dt.date
    trend_df = df.groupby(['Date', 'Entity']).size().unstack(fill_value=0)
    
    # Sort by date
    trend_df = trend_df.sort_index()
    
    # Create chart data
    chart_data = CategoryChartData()
    chart_data.categories = [str(date) for date in trend_df.index]
    
    entities = trend_df.columns.tolist()
    for entity in entities:
        chart_data.add_series(entity, trend_df[entity].tolist())
    
    # Generate colors for entities (cycling through severity colors)
    colors_list = [COLORS['critical'], COLORS['high'], COLORS['medium'], 
                   COLORS['low'], COLORS['accent'], COLORS['wip']]
    series_colors = [colors_list[i % len(colors_list)] for i in range(len(entities))]
    
    # Add chart
    create_line_chart(
        slide, chart_data,
        "Daily Alert Trend by Entity",
        Inches(0.5), Inches(1.2), Inches(9), Inches(3.8),
        series_colors
    )
    
    # Add note
    note_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(0.5))
    note_frame = note_box.text_frame
    note_frame.text = "Note: This chart is fully editable. Right-click the chart to edit data and formatting."
    note_frame.paragraphs[0].font.size = Pt(11)
    note_frame.paragraphs[0].font.italic = True
    note_frame.paragraphs[0].font.color.rgb = COLORS['text']


def generate_incident_summary_slide(prs, df):
    """Slide 5: Weekly Incident Summary Table"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    # Add slide title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "Weekly Incident Summary"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLORS['accent']
    
    # Prepare comprehensive summary data
    summary_data = []
    
    for entity in df['Entity'].unique():
        entity_df = df[df['Entity'] == entity]
        
        for severity in SEVERITY_ORDER:
            severity_df = entity_df[entity_df['Alert Severity'] == severity]
            
            if len(severity_df) > 0:
                incident_count = len(severity_df)
                tp_count = len(severity_df[severity_df['Alert Efficiency'] == 'True Positive'])
                fp_count = len(severity_df[severity_df['Alert Efficiency'] == 'False Positive'])
                
                # Determine category based on severity
                if severity == 'Critical':
                    category = 'Security Incident'
                elif severity == 'High':
                    category = 'Security Event'
                else:
                    category = 'Alert'
                
                summary_data.append({
                    'Entity': entity,
                    'Severity': severity,
                    'Incidents': incident_count,
                    'TP': tp_count,
                    'FP': fp_count,
                    'Category': category,
                    'Remarks': 'Under Investigation' if tp_count > 0 else 'Resolved'
                })
    
    # Create table
    table_data = [['Entity', 'Severity', 'Count', 'TP', 'FP', 'Category', 'Remarks']]
    
    for item in summary_data:
        row = [
            item['Entity'],
            item['Severity'],
            item['Incidents'],
            item['TP'],
            item['FP'],
            item['Category'],
            item['Remarks']
        ]
        table_data.append(row)
    
    # Add table with smaller font for more data
    table_placeholder = slide.shapes.add_table(
        len(table_data), 7, Inches(0.3), Inches(1.0), Inches(9.4), Inches(6)
    )
    table = table_placeholder.table
    
    # Populate and format table
    for i, row_data in enumerate(table_data):
        for j, cell_value in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = str(cell_value)
            
            # Format cell
            cell.text_frame.paragraphs[0].font.size = Pt(9)
            
            # Header row formatting
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['accent']
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            else:
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                # Color code severity column
                if j == 1:  # Severity column
                    severity_value = str(cell_value)
                    if severity_value == 'Critical':
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(255, 230, 230)
                    elif severity_value == 'High':
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(255, 240, 220)


def generate_soc_report(excel_file_path, output_dir=None):
    """
    Main function to generate SOC report
    
    Args:
        excel_file_path: Path to Excel file with EDR alert data
        output_dir: Optional output directory (defaults to same as Excel file)
    
    Returns:
        Path to generated PowerPoint file
    """
    print("Starting SOC Report Generation...")
    
    # Read data
    print(f"Reading data from: {excel_file_path}")
    df = read_excel_data(excel_file_path)
    print(f"Successfully read {len(df)} records")
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Get date range for report
    min_date = df['Alert Trend'].min().strftime('%Y-%m-%d')
    max_date = df['Alert Trend'].max().strftime('%Y-%m-%d')
    
    # Add title slide
    print("Creating title slide...")
    add_title_slide(
        prs,
        "EDR Weekly Incident Report",
        f"Security Operations Center\n{min_date} to {max_date}"
    )
    
    # Generate content slides
    print("Generating Alert Severity Distribution slide...")
    generate_severity_slide(prs, df)
    
    print("Generating Alert Status Overview slide...")
    generate_status_slide(prs, df)
    
    print("Generating Alert Efficiency Analysis slide...")
    generate_efficiency_slide(prs, df)
    
    print("Generating Alert Trend Analysis slide...")
    generate_trend_slide(prs, df)
    
    print("Generating Weekly Incident Summary slide...")
    generate_incident_summary_slide(prs, df)
    
    # Save presentation
    if output_dir is None:
        output_dir = os.path.dirname(excel_file_path)
    
    current_date = datetime.now().strftime('%Y-%m-%d')
    output_file = os.path.join(output_dir, f"EDR-Weekly-Incident-Report_{current_date}.pptx")
    
    prs.save(output_file)
    print(f"\n[SUCCESS] Report generated successfully!")
    print(f"[SUCCESS] Output file: {output_file}")
    
    return output_file


if __name__ == "__main__":
    # Example usage
    print("=" * 60)
    print("SOC Weekly Report Generator")
    print("=" * 60)
    
    # Get Excel file path from user
    excel_file = input("\nEnter the path to your Excel file: ").strip().strip('"\'')
    
    if not os.path.exists(excel_file):
        print(f"Error: File not found: {excel_file}")
    else:
        try:
            output_path = generate_soc_report(excel_file)
            print(f"\n{'=' * 60}")
            print("Report generation completed!")
            print(f"{'=' * 60}")
        except Exception as e:
            print(f"\nError generating report: {str(e)}")
            import traceback
            traceback.print_exc()
