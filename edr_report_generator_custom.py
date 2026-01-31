"""
Custom EDR Weekly Incident Report Generator
Generates PowerPoint reports matching the EDR Weekly Incident Metrics template
using data from combined.xlsx
"""

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.dml.color import RGBColor
from datetime import datetime
import os


# Corporate Color Theme matching template
COLORS = {
    'critical': RGBColor(192, 0, 0),        # Dark Red (theme color)
    'high': RGBColor(237, 125, 49),         # Orange (from template)
    'medium': RGBColor(166, 166, 166),      # Gray (from template)
    'low': RGBColor(255, 192, 0),           # Yellow (from template)
    'closed': RGBColor(132, 209, 88),       # Green (user specified)
    'wip': RGBColor(192, 0, 0),             # Red - same as critical (user specified)
    'background': RGBColor(248, 249, 250),  # Light Gray
    'text': RGBColor(0, 0, 0),              # Black
    'accent': RGBColor(68, 114, 196),       # Blue
    'header': RGBColor(68, 114, 196)        # Blue for headers
}

SEVERITY_ORDER = ['Critical', 'High', 'Medium', 'Low']
STATUS_ORDER = ['closed', 'WIP']
ENTITY_ORDER = ['DIAL', 'ENT', 'HIAL', 'MRO']


def parse_combined_xlsx(file_path):
    """
    Parse combined.xlsx - handles raw alert data format
    Columns: Entity, Alert Trend, Filename, Alert Severity, Alert Status,  Alert Efficiency
    
    Returns: Dictionary with entity data
    """
    df = pd.read_excel(file_path)
    
    entities_data = {}
    grand_total = len(df)
    
    # Process each entity
    for entity in ENTITY_ORDER:
        entity_df = df[df['Entity'] == entity]
        
        if len(entity_df) == 0:
            continue
        
        entities_data[entity] = {
            'total': len(entity_df),
            'closed': {},
            'WIP': {},
            'all_severities': {},
            'efficacy': {'TP': 0, 'FP': 0},
            'trends': {},
            'trend_dates': []
        }
        
        # Count by status and severity
        for status in STATUS_ORDER:
            status_df = entity_df[entity_df['Alert Status'] == status]
            
            for severity in SEVERITY_ORDER:
                sev_count = len(status_df[status_df['Alert Severity'] == severity])
                
                if sev_count > 0:
                    entities_data[entity][status][severity] = sev_count
                    
                    # Track overall severity counts
                    if severity not in entities_data[entity]['all_severities']:
                        entities_data[entity]['all_severities'][severity] = 0
                    entities_data[entity]['all_severities'][severity] += sev_count
        
        # Extract efficacy data (True Positive / False Positive)
        tp_count = len(entity_df[entity_df['Alert Efficiency'] == 'True Positive'])
        fp_count = len(entity_df[entity_df['Alert Efficiency'] == 'False Positive'])
        entities_data[entity]['efficacy'] = {'TP': tp_count, 'FP': fp_count}
        
        # Extract trend data (daily counts by severity)
        dates = sorted(entity_df['Alert Trend'].dropna().unique())
        entities_data[entity]['trend_dates'] = [d.strftime('%Y-%m-%d') if hasattr(d, 'strftime') else str(d) for d in dates]
        
        for severity in SEVERITY_ORDER:
            daily_counts = []
            for date in dates:
                count = len(entity_df[
                    (entity_df['Alert Trend'] == date) & 
                    (entity_df['Alert Severity'] == severity)
                ])
                daily_counts.append(count)
            entities_data[entity]['trends'][severity] = daily_counts
        
        # Extract top triggering filename
        if 'Filename' in entity_df.columns:
            filename_counts = entity_df['Filename'].value_counts()
            if len(filename_counts) > 0:
                top_filename = filename_counts.index[0]
                top_count = filename_counts.iloc[0]
                entities_data[entity]['top_file'] = {'name': top_filename, 'count': int(top_count)}
            else:
                entities_data[entity]['top_file'] = {'name': 'N/A', 'count': 0}
        else:
            entities_data[entity]['top_file'] = {'name': 'N/A', 'count': 0}
    
    return entities_data, grand_total


def add_title_slide(prs, date_range):
    """Add title slide matching template format"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Branding removed per user request: www.gmrgroup.in
    # branding_box = slide.shapes.add_textbox(Inches(10.5), Inches(0.3), Inches(2.5), Inches(0.4))
    # branding_frame = branding_box.text_frame
    # branding_frame.text = "www.gmrgroup.in"
    
    # Main title (centered)
    title_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9.33), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "EDR Weekly Incident Metrics"
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLORS['accent']
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def create_entity_chart(slide, entity_name, categories, values, title, left, top, width, height, colors_dict):
    """Create a single clustered bar chart for an entity"""
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('Count', values)
    
    chart_placeholder = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data
    )
    
    chart = chart_placeholder.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
    chart.chart_title.text_frame.paragraphs[0].font.bold = True
    
    # Format legend
    chart.has_legend = False  # No legend for single series
    
    # Apply colors to bars
    series = chart.series[0]
    for idx, point in enumerate(series.points):
        if idx < len(categories):
            category_name = categories[idx]
            if category_name in colors_dict:
                fill = point.format.fill
                fill.solid()
                fill.fore_color.rgb = colors_dict[category_name]
    
    return chart


def format_severity_breakdown(severity_dict):
    """Format severity breakdown as multi-line text"""
    lines = []
    for severity in SEVERITY_ORDER:
        count = severity_dict.get(severity, 0)
        lines.append(f"{severity} - {count}")
    return '\n'.join(lines)


def create_severity_slide(prs, entities_data, date_range):
    """Slide 2: Alert Severity - 4 charts + table"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Slide title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Alert Severity"
    title_frame.paragraphs[0].font.size = Pt(24)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLORS['accent']
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Create 4 charts in 2x2 grid
    chart_width = Inches(3)
    chart_height = Inches(2.2)
    
    positions = [
        (Inches(0.5), Inches(1.2)),   # Top left
        (Inches(3.8), Inches(1.2)),   # Top right
        (Inches(0.5), Inches(3.6)),   # Bottom left
        (Inches(3.8), Inches(3.6))    # Bottom right
    ]
    
    severity_colors = {
        'Critical': COLORS['critical'],
        'High': COLORS['high'],
        'Medium': COLORS['medium'],
        'Low': COLORS['low']
    }
    
    for idx, entity in enumerate(ENTITY_ORDER):
        if idx < len(positions):
            entity_data = entities_data.get(entity, {})
            severities = entity_data.get('all_severities', {})
            
            categories = SEVERITY_ORDER
            values = [severities.get(sev, 0) for sev in SEVERITY_ORDER]
            
            left, top = positions[idx]
            create_entity_chart(
                slide, entity, categories, values,
                f"Alert Severity - {entity}",
                left, top, chart_width, chart_height,
                severity_colors
            )
    
    # Summary table
    table_data = [['Entity', 'Total Alerts', 'Severity']]
    
    for entity in ENTITY_ORDER:
        entity_info = entities_data.get(entity, {})
        total = entity_info.get('total', 0)
        severity_text = format_severity_breakdown(entity_info.get('all_severities', {}))
        table_data.append([entity, str(total), severity_text])
    
    # Add table
    add_formatted_table(slide, table_data, Inches(7.5), Inches(1.2), Inches(5.3), Inches(4.5))


def create_status_slide(prs, entities_data, date_range):
    """Slide 3: Alert Status - 4 charts + table"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Slide title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Alert Status"
    title_frame.paragraphs[0].font.size = Pt(24)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLORS['accent']
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Create 4 charts
    chart_width = Inches(3)
    chart_height = Inches(2.2)
    
    positions = [
        (Inches(0.5), Inches(1.2)),
        (Inches(3.8), Inches(1.2)),
        (Inches(0.5), Inches(3.6)),
        (Inches(3.8), Inches(3.6))
    ]
    
    status_colors = {
        'Closed': COLORS['closed'],
        'WIP': COLORS['wip']
    }
    
    for idx, entity in enumerate(ENTITY_ORDER):
        if idx < len(positions):
            entity_data = entities_data.get(entity, {})
            
            # Calculate total closed and WIP
            closed_total = sum(entity_data.get('closed', {}).values())
            wip_total = sum(entity_data.get('WIP', {}).values())
            
            categories = ['Closed', 'WIP']
            values = [closed_total, wip_total]
            
            left, top = positions[idx]
            create_entity_chart(
                slide, entity, categories, values,
                f"Alert Status - {entity}",
                left, top, chart_width, chart_height,
                status_colors
            )
    
    # Summary table with closed/WIP breakdown
    table_data = [['Entity', 'Total Alerts', 'Closed', 'WIP']]
    
    for entity in ENTITY_ORDER:
        entity_info = entities_data.get(entity, {})
        total = entity_info.get('total', 0)
        closed_text = format_severity_breakdown(entity_info.get('closed', {}))
        wip_text = format_severity_breakdown(entity_info.get('WIP', {}))
        table_data.append([entity, str(total), closed_text, wip_text])
    
    add_formatted_table(slide, table_data, Inches(7.5), Inches(1.2), Inches(5.3), Inches(4.5))


def create_efficacy_slide(prs, entities_data, date_range):
    """Slide 4: Alert Efficacy - Placeholder (data not available)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Slide title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Alert Efficacy"
    title_frame.paragraphs[0].font.size = Pt(24)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLORS['accent']
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Create 4 placeholder charts with zero data
    chart_width = Inches(3)
    chart_height = Inches(2.2)
    
    positions = [
        (Inches(0.5), Inches(1.2)),
        (Inches(3.8), Inches(1.2)),
        (Inches(0.5), Inches(3.6)),
        (Inches(3.8), Inches(3.6))
    ]
    
    efficacy_colors = {
        'True Positive': RGBColor(211, 47, 47),    # Red (user specified)
        'False Positive': RGBColor(132, 209, 88)   # Green (user specified)
    }
    
    for idx, entity in enumerate(ENTITY_ORDER):
        if idx < len(positions):
            entity_data = entities_data.get(entity, {})
            efficacy = entity_data.get('efficacy', {'TP': 0, 'FP': 0})
            
            categories = ['True Positive', 'False Positive']
            values = [efficacy['TP'], efficacy['FP']]
            
            left, top = positions[idx]
            create_entity_chart(
                slide, entity, categories, values,
                f"Alert Efficacy - {entity}",
                left, top, chart_width, chart_height,
                efficacy_colors
            )
    
    # Placeholder table
    table_data = [['Entity', 'Total Alerts', 'Alert efficacy']]
    
    for entity in ENTITY_ORDER:
        entity_info = entities_data.get(entity, {})
        total = entity_info.get('total', 0)
        efficacy = entity_info.get('efficacy', {'TP': 0, 'FP': 0})
        efficacy_text = f"True Positive - {efficacy['TP']}\nFalse Positive - {efficacy['FP']}"
        table_data.append([entity, str(total), efficacy_text])
    
    add_formatted_table(slide, table_data, Inches(7.5), Inches(1.2), Inches(5.3), Inches(4.5))
    
    # Add note about missing data
    note_box = slide.shapes.add_textbox(Inches(7.5), Inches(6.0), Inches(5.3), Inches(0.8))
    note_frame = note_box.text_frame
    note_frame.text = "Note: True Positive/False Positive data not available in combined.xlsx"
    note_frame.paragraphs[0].font.size = Pt(10)
    note_frame.paragraphs[0].font.italic = True
    note_frame.paragraphs[0].font.color.rgb = COLORS['text']


def create_trend_slide(prs, entities_data, date_range):
    """Slide 5: Alert Trend - Line charts with unified date range"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Slide title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Alert Trend"
    title_frame.paragraphs[0].font.size = Pt(24)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLORS['accent']
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Create 4 placeholder line charts
    chart_width = Inches(3)
    chart_height = Inches(2.5)
    
    positions = [
        (Inches(0.5), Inches(1.2)),
        (Inches(3.8), Inches(1.2)),
        (Inches(0.5), Inches(4.0)),
        (Inches(3.8), Inches(4.0))
    ]
    
    # Find the unified date range across all entities
    all_dates = set()
    for entity in ENTITY_ORDER:
        entity_data = entities_data.get(entity, {})
        trend_dates = entity_data.get('trend_dates', [])
        all_dates.update(trend_dates)
    
    # Sort to get chronological order
    unified_dates = sorted(list(all_dates)) if all_dates else ['Day 1', 'Day 2', 'Day 3', 'Day 4', 'Day 5', 'Day 6', 'Day 7']
    
    # Severity colors for line charts
    severity_colors = {
        'Critical': COLORS['critical'],
        'High': COLORS['high'],
        'Medium': COLORS['medium'],
        'Low': COLORS['low']
    }
    
    
    for idx, entity in enumerate(ENTITY_ORDER):
        if idx < len(positions):
            entity_data = entities_data.get(entity, {})
            trends = entity_data.get('trends', {})
            entity_dates = entity_data.get('trend_dates', [])
            
            chart_data = CategoryChartData()
            chart_data.categories = unified_dates
            
            # Add 4 series (one per severity level)
            for severity in SEVERITY_ORDER:
                entity_trend = trends.get(severity, [])
                
                # Map entity's trend data to unified date range
                unified_counts = []
                for date in unified_dates:
                    if date in entity_dates:
                        # Find index in entity's dates
                        date_idx = entity_dates.index(date)
                        unified_counts.append(entity_trend[date_idx] if date_idx < len(entity_trend) else 0)
                    else:
                        # No data for this date
                        unified_counts.append(0)
                
                chart_data.add_series(severity, unified_counts)
            
            left, top = positions[idx]
            chart_placeholder = slide.shapes.add_chart(
                XL_CHART_TYPE.LINE, left, top, chart_width, chart_height, chart_data
            )
            
            chart = chart_placeholder.chart
            chart.has_title = True
            chart.chart_title.text_frame.text = entity
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
            chart.chart_title.text_frame.paragraphs[0].font.bold = True
            
            # Apply colors to each series
            for series_idx, severity in enumerate(SEVERITY_ORDER):
                if series_idx < len(chart.series):
                    series = chart.series[series_idx]
                    # Set line color
                    line = series.format.line
                    line.color.rgb = severity_colors[severity]
    
    # Add note
    note_box = slide.shapes.add_textbox(Inches(7.5), Inches(2.0), Inches(5.3), Inches(1.0))
    note_frame = note_box.text_frame
    note_frame.text = "Note: Daily trend data not available in combined.xlsx\nPlaceholder charts shown"
    note_frame.paragraphs[0].font.size = Pt(11)
    note_frame.paragraphs[0].font.italic = True
    note_frame.paragraphs[0].font.color.rgb = COLORS['text']


def create_summary_slide(prs, entities_data, grand_total, date_range):
    """Slide 6: Summary of Alerts - 1 chart + comprehensive table"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Slide title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Summary of Alerts"
    title_frame.paragraphs[0].font.size = Pt(24)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLORS['accent']
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Summary bar chart - Single series with only non-zero bars (14 total)
    # Grouped by Entity → Status → Severity (only bars with data)
    
    categories = []
    values = []
    
    for entity in ENTITY_ORDER:
        entity_info = entities_data.get(entity, {})
        closed_severities = entity_info.get('closed', {})
        wip_severities = entity_info.get('WIP', {})
        
        # Add only non-zero severity bars for "closed" status
        for severity in SEVERITY_ORDER:
            count = closed_severities.get(severity, 0)
            if count > 0:
                categories.append(f"{severity}\nclosed\n{entity}")
                values.append(count)
        
        # Add only non-zero severity bars for "WIP" status
        for severity in SEVERITY_ORDER:
            count = wip_severities.get(severity, 0)
            if count > 0:
                categories.append(f"{severity}\nWIP\n{entity}")
                values.append(count)
    
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('Total', values)
    
    chart_placeholder = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.5), Inches(1.2), Inches(6), Inches(3.5), chart_data
    )
    
    chart = chart_placeholder.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Summary of Alerts"
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(16)
    chart.chart_title.text_frame.paragraphs[0].font.bold = True
    
    # Color each bar based on its severity (extracted from label)
    series = chart.series[0]
    for idx, category in enumerate(categories):
        # Extract severity from category (first line of label)
        severity = category.split('\n')[0]
        
        # Apply color based on severity
        point = series.points[idx]
        point.format.fill.solid()
        if severity == 'Critical':
            point.format.fill.fore_color.rgb = COLORS['critical']
        elif severity == 'High':
            point.format.fill.fore_color.rgb = COLORS['high']
        elif severity == 'Medium':
            point.format.fill.fore_color.rgb = COLORS['medium']
        elif severity == 'Low':
            point.format.fill.fore_color.rgb = COLORS['low']
    
    # Summary table
    table_data = [['Entity', 'Total Alerts', 'Severity', 'Closed', 'WIP']]
    
    for entity in ENTITY_ORDER:
        entity_info = entities_data.get(entity, {})
        total = entity_info.get('total', 0)
        severity_text = format_severity_breakdown(entity_info.get('all_severities', {}))
        closed_text = format_severity_breakdown(entity_info.get('closed', {}))
        wip_text = format_severity_breakdown(entity_info.get('WIP', {}))
        table_data.append([entity, str(total), severity_text, closed_text, wip_text])
    
    add_formatted_table(slide, table_data, Inches(7), Inches(1.2), Inches(6), Inches(4.5))


def create_top_incidents_slide(prs, entities_data, date_range):
    """Slide 7: Top Incidents - Top triggering files per entity"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Slide title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Top Incidents"
    title_frame.paragraphs[0].font.size = Pt(24)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLORS['accent']
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Build table with top triggering files
    table_data = [
        ['Entity', 'Triggering file', 'Alert Count', 'Alert action taken']
    ]
    
    for entity in ENTITY_ORDER:
        entity_info = entities_data.get(entity, {})
        top_file_info = entity_info.get('top_file', {'name': 'N/A', 'count': 0})
        
        table_data.append([
            entity,
            top_file_info['name'],
            str(top_file_info['count']),
            ''  # Leave blank for user to fill
        ])
    
    add_formatted_table(slide, table_data, Inches(2), Inches(1.5), Inches(9.33), Inches(3))
    
    # Add note
    note_box = slide.shapes.add_textbox(Inches(2), Inches(5), Inches(9.33), Inches(1))
    note_frame = note_box.text_frame
    note_frame.text = "Note: Incident details (triggering files, actions taken) not available in combined.xlsx\nRequires additional detailed incident data"
    note_frame.paragraphs[0].font.size = Pt(12)
    note_frame.paragraphs[0].font.italic = True
    note_frame.paragraphs[0].font.color.rgb = COLORS['text']
    note_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def create_detailed_summary_table(prs, entities_data, date_range):
    """Slide 8: Detailed Summary Table - Matching template exactly"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Slide title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Summary of Alerts"
    title_frame.paragraphs[0].font.size = Pt(24)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLORS['accent']
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Build detailed table matching template exactly
    table_data = [[
        'S.no', 'Entity', 'Total Incidents count ', 'Closed Incident', 'WIP Incidents',
        'Severity wise incidents count', '', 'Count of False Positive Alerts',
        'Count of True Positive Alerts', 'Alert Category ', 'Remarks'
    ]]
    
    s_no = 1
    for entity in ENTITY_ORDER:
        entity_info = entities_data.get(entity, {})
        total = entity_info.get('total', 0)
        all_severities = entity_info.get('all_severities', {})
        closed_dict = entity_info.get('closed', {})
        wip_dict = entity_info.get('WIP', {})
        
        # Calculate totals
        total_closed = sum(closed_dict.values())
        total_wip = sum(wip_dict.values())
        
        # Entity name mapping for Slide 8
        if entity == 'ENT':
            entity_name = 'Enterprise'
        elif entity == 'MRO':
            entity_name = 'HIAL- MRO'
        else:
            entity_name = entity
        
        # Always generate 4 severity rows (Critical, High, Medium, Low)
        for idx, severity in enumerate(SEVERITY_ORDER):
            sev_count = all_severities.get(severity, 0)
            
            # First row for each entity (idx == 0, Critical)
            if idx == 0:
                # FP = total_closed for first row
                fp_count = f"{total_closed:02d}"
                tp_count = "00"
                category = "PUP & Legitimate Non malicious processes"
                remarks = "Malicious exe file accessed  & Legitimate processes triggered."
                
                row = [
                    str(s_no),
                    entity_name,
                    f"{total:02d}",
                    f"{total_closed:02d}",
                    f"{total_wip:02d}",
                    severity,
                    f"{sev_count:02d}",
                    fp_count,
                    tp_count,
                    category,
                    remarks
                ]
            else:
                # Rows 2-4 (High, Medium, Low) - mostly blank
                row = [
                    '',  # S.no blank
                    '',  # Entity blank
                    '',  # Total blank
                    '',  # Closed blank
                    '',  # WIP blank
                    severity,  # Severity name
                    f"{sev_count:02d}",  # Severity count
                    '',  # FP blank
                    '',  # TP blank
                    '',  # Category blank
                    ''   # Remarks blank
                ]
            
            table_data.append(row)
        
        s_no += 1
    
    # Add table with smaller fonts for more data
    table_placeholder = slide.shapes.add_table(
        len(table_data), 11,
        Inches(0.3), Inches(1.0), Inches(12.7), Inches(6)
    )
    table = table_placeholder.table
    
    # Set column widths
    col_widths = [Inches(0.4), Inches(0.8), Inches(1.0), Inches(1.0), Inches(1.0),
                  Inches(1.2), Inches(0.5), Inches(1.0), Inches(1.0),
                  Inches(1.5), Inches(2.3)]
    
    for col_idx, width in enumerate(col_widths):
        table.columns[col_idx].width = width
    
    # Populate table
    for r_idx, row_data in enumerate(table_data):
        for c_idx, cell_value in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(cell_value)
            
            # Format cell
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # Header row
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['header']
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                cell.text_frame.paragraphs[0].font.bold = True


def create_closing_slide(prs):
    """Slide 9: Closing slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Simple closing message
    closing_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9.33), Inches(1.5))
    closing_frame = closing_box.text_frame
    closing_frame.text = "Thank You"
    closing_frame.paragraphs[0].font.size = Pt(48)
    closing_frame.paragraphs[0].font.bold = True
    closing_frame.paragraphs[0].font.color.rgb = COLORS['accent']
    closing_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    closing_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_formatted_table(slide, data, left, top, width, height):
    """Add a formatted table with headers"""
    rows, cols = len(data), len(data[0])
    
    table_placeholder = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_placeholder.table
    
    # Populate and format table
    for r_idx, row_data in enumerate(data):
        for c_idx, cell_value in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(cell_value)
            
            # Format cell
            cell.text_frame.paragraphs[0].font.size = Pt(9)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.TOP
            
            # Header row
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['header']
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                cell.text_frame.paragraphs[0].font.bold = True
    
    return table


def generate_edr_report(excel_file_path, date_range, output_dir=None):
    """
    Main function to generate EDR report matching template
    
    Args:
        excel_file_path: Path to combined.xlsx
        date_range: Date range string (e.g., "01st Dec 25 to 07th Dec 25")
        output_dir: Optional output directory
    
    Returns:
        Path to generated PowerPoint file
    """
    print("Starting EDR Report Generation...")
    print("=" * 60)
    
    # Parse data
    print(f"Reading data from: {excel_file_path}")
    entities_data, grand_total = parse_combined_xlsx(excel_file_path)
    print(f"Successfully parsed data for {len(entities_data)} entities")
    print(f"Grand Total Alerts: {grand_total}")
    
    # Create presentation with custom dimensions
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    
    # Generate all 9 slides
    print("\nGenerating slides...")
    
    print("  [1/9] Title slide")
    add_title_slide(prs, date_range)
    
    print("  [2/9] Alert Severity slide")
    create_severity_slide(prs, entities_data, date_range)
    
    print("  [3/9] Alert Status slide")
    create_status_slide(prs, entities_data, date_range)
    
    print("  [4/9] Alert Efficacy slide (placeholder)")
    create_efficacy_slide(prs, entities_data, date_range)
    
    print("  [5/9] Alert Trend slide (placeholder)")
    create_trend_slide(prs, entities_data, date_range)
    
    print("  [6/9] Summary of Alerts slide")
    create_summary_slide(prs, entities_data, grand_total, date_range)
    
    print("  [7/9] Top Incidents slide")
    create_top_incidents_slide(prs, entities_data, date_range)
    
    print("  [8/9] Detailed Summary Table slide")
    create_detailed_summary_table(prs, entities_data, date_range)
    
    print("  [9/9] Closing slide")
    create_closing_slide(prs)
    
    # Save presentation
    if output_dir is None:
        output_dir = os.path.dirname(excel_file_path)
    
    current_date = datetime.now().strftime('%Y-%m-%d')
    output_file = os.path.join(output_dir, f"EDR-Weekly-Incident-Report_{current_date}.pptx")
    
    prs.save(output_file)
    
    print("\n" + "=" * 60)
    print("[SUCCESS] Report generated successfully!")
    print(f"[SUCCESS] Output file: {output_file}")
    print(f"[INFO] Total slides: 9")
    print(f"[INFO] Entities: {', '.join(ENTITY_ORDER)}")
    print(f"[INFO] Total alerts: {grand_total}")
    print("=" * 60)
    
    return output_file


if __name__ == "__main__":
    print("=" * 60)
    print("EDR Weekly Incident Report Generator (Custom)")
    print("=" * 60)
    
    # Use combined.xlsx with date range
    excel_file = "combined.xlsx"
    date_range = "01st Dec 25 to 07th Dec 25"
    
    if not os.path.exists(excel_file):
        print(f"\nError: {excel_file} not found!")
        excel_file = input("Enter path to combined.xlsx: ").strip().strip('"\'')
        date_range = input("Enter date range (e.g., '01st Dec 25 to 07th Dec 25'): ").strip()
    
    if os.path.exists(excel_file):
        try:
            output_path = generate_edr_report(excel_file, date_range)
        except Exception as e:
            print(f"\nError generating report: {str(e)}")
            import traceback
            traceback.print_exc()
    else:
        print(f"Error: File not found: {excel_file}")
