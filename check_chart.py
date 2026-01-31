"""Quick check of the generated chart"""
from pptx import Presentation

prs = Presentation("EDR-Weekly-Incident-Report_2026-01-31.pptx")
slide6 = prs.slides[5]
chart = [s for s in slide6.shapes if hasattr(s, 'has_chart') and s.has_chart][0].chart

print("Chart Structure:")
print(f"Number of series: {len(chart.series)}")
print(f"Number of bars: {len(list(chart.series[0].values))}")

categories = list(chart.plots[0].categories)
values = list(chart.series[0].values)

print(f"\nAll {len(categories)} bars:")
for i, (cat, val) in enumerate(zip(categories, values)):
    print(f"{i+1}. Label: '{cat}' = {val}")
