"""
Detailed analysis of template Slide 6 chart vs table
"""
from pptx import Presentation

template = Presentation("EDR Weekly Incident Metrics (01st Dec 25 to 07th Dec 25).pptx")
slide6 = template.slides[5]

print("=" * 80)
print("TEMPLATE SLIDE 6 ANALYSIS")
print("=" * 80)

# Get chart
charts = [s for s in slide6.shapes if hasattr(s, 'has_chart') and s.has_chart]
if charts:
    chart = charts[0].chart
    print("\nCHART:")
    print(f"  Title: {chart.chart_title.text_frame.text if chart.has_title else 'No title'}")
    print(f"  Series count: {len(chart.series)}")
    
    for idx, series in enumerate(chart.series):
        values = list(series.values)
        categories = list(chart.plots[0].categories)
        print(f"\n  Series {idx+1}: '{series.name}'")
        print(f"    Total bars: {len(values)}")
        print(f"    Values: {values}")
        print(f"    Sum: {sum(values)}")
        print(f"\n    Categories: {categories}")
        
        # Try to decode the pattern
        print(f"\n    Bar-by-bar breakdown:")
        for i, (cat, val) in enumerate(zip(categories, values)):
            print(f"      {i+1}. {cat}: {val}")

# Get table
tables = [s for s in slide6.shapes if hasattr(s, 'has_table') and s.has_table]
if tables:
    table = tables[0].table
    print(f"\n\nTABLE: {len(table.rows)} rows × {len(table.columns)} cols")
    for r in range(min(6, len(table.rows))):
        row_data = [table.cell(r, c).text for c in range(len(table.columns))]
        print(f"  Row {r}: {row_data}")

print("\n" + "=" * 80)
