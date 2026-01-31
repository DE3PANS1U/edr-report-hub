"""
Extract exact RGB colors from template chart bars
"""
from pptx import Presentation

template = Presentation("EDR Weekly Incident Metrics (01st Dec 25 to 07th Dec 25).pptx")
slide6 = template.slides[5]

charts = [s for s in slide6.shapes if hasattr(s, 'has_chart') and s.has_chart]
chart = charts[0].chart

print("Template Chart Bar Colors:")
print("=" * 60)

categories = list(chart.plots[0].categories)
values = list(chart.series[0].values)

# Try to get colors from chart points
series = chart.series[0]
for idx in range(min(14, len(series.points))):
    point = series.points[idx]
    category = categories[idx]
    value = values[idx]
    
    try:
        # Try to get fill color
        rgb = point.format.fill.fore_color.rgb
        print(f"{idx+1}. {category:10} value={value:4.0f} RGB=({rgb[0]:3}, {rgb[1]:3}, {rgb[2]:3}) = RGBColor({rgb[0]}, {rgb[1]}, {rgb[2]})")
    except Exception as e:
        print(f"{idx+1}. {category:10} value={value:4.0f} Error getting color")

print("\n" + "=" * 60)
