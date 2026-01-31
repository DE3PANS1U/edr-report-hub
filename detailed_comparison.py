"""
Detailed comparison of template vs generated PowerPoint
"""
from pptx import Presentation
import os

template_file = "EDR Weekly Incident Metrics (01st Dec 25 to 07th Dec 25).pptx"
generated_file = "EDR-Weekly-Incident-Report_2026-01-31.pptx"

print("=" * 80)
print("DETAILED TEMPLATE ANALYSIS")
print("=" * 80)

if os.path.exists(template_file):
    template_prs = Presentation(template_file)
    
    for slide_idx, slide in enumerate(template_prs.slides, 1):
        print(f"\n{'=' * 80}")
        print(f"TEMPLATE SLIDE {slide_idx}")
        print(f"{'=' * 80}")
        
        # Get all text
        all_text = []
        chart_count = 0
        table_count = 0
        
        for shape in slide.shapes:
            # Text
            if hasattr(shape, "text") and shape.text.strip():
                all_text.append(shape.text.strip())
            
            # Charts
            if hasattr(shape, "has_chart") and shape.has_chart:
                chart_count += 1
                chart = shape.chart
                print(f"\nChart {chart_count}:")
                print(f"  Type: {chart.chart_type}")
                if chart.has_title:
                    print(f"  Title: {chart.chart_title.text_frame.text}")
                
                # Get chart data
                try:
                    for series_idx, series in enumerate(chart.series):
                        print(f"  Series {series_idx + 1}: {series.name}")
                        # Try to get values
                        try:
                            print(f"    Values: {list(series.values)[:10]}")
                        except:
                            pass
                except Exception as e:
                    print(f"  (Could not extract chart data: {e})")
            
            # Tables
            if hasattr(shape, "has_table") and shape.has_table:
                table_count += 1
                table = shape.table
                print(f"\nTable {table_count}: {len(table.rows)} rows x {len(table.columns)} columns")
                
                # Print all table content
                for r_idx in range(len(table.rows)):
                    row_data = []
                    for c_idx in range(len(table.columns)):
                        cell_text = table.cell(r_idx, c_idx).text.strip()
                        row_data.append(cell_text[:50])
                    print(f"  Row {r_idx}: {row_data}")
        
        print(f"\nText boxes: {len(all_text)}")
        for idx, text in enumerate(all_text[:5], 1):
            print(f"  {idx}. {text[:100]}")

print("\n\n" + "=" * 80)
print("DETAILED GENERATED FILE ANALYSIS")
print("=" * 80)

if os.path.exists(generated_file):
    generated_prs = Presentation(generated_file)
    
    for slide_idx, slide in enumerate(generated_prs.slides, 1):
        print(f"\n{'=' * 80}")
        print(f"GENERATED SLIDE {slide_idx}")
        print(f"{'=' * 80}")
        
        # Get all text
        all_text = []
        chart_count = 0
        table_count = 0
        
        for shape in slide.shapes:
            # Text
            if hasattr(shape, "text") and shape.text.strip():
                all_text.append(shape.text.strip())
            
            # Charts
            if hasattr(shape, "has_chart") and shape.has_chart:
                chart_count += 1
                chart = shape.chart
                print(f"\nChart {chart_count}:")
                print(f"  Type: {chart.chart_type}")
                if chart.has_title:
                    print(f"  Title: {chart.chart_title.text_frame.text}")
            
            # Tables
            if hasattr(shape, "has_table") and shape.has_table:
                table_count += 1
                table = shape.table
                print(f"\nTable {table_count}: {len(table.rows)} rows x {len(table.columns)} columns")
                
                # Print all table content
                for r_idx in range(min(10, len(table.rows))):
                    row_data = []
                    for c_idx in range(len(table.columns)):
                        cell_text = table.cell(r_idx, c_idx).text.strip()
                        row_data.append(cell_text[:50])
                    print(f"  Row {r_idx}: {row_data}")
        
        print(f"\nText boxes: {len(all_text)}")
        for idx, text in enumerate(all_text[:5], 1):
            print(f"  {idx}. {text[:100]}")

print("\n\n" + "=" * 80)
print("COMPARISON SUMMARY")
print("=" * 80)
print(f"Template slides: {len(template_prs.slides)}")
print(f"Generated slides: {len(generated_prs.slides)}")
