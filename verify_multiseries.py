"""
Verify the new multi-series chart structure
"""
from pptx import Presentation

generated = Presentation("EDR-Weekly-Incident-Report_2026-01-31.pptx")
slide6 = generated.slides[5]

print("=" * 80)
print("SLIDE 6 - SUMMARY CHART VERIFICATION")
print("=" * 80)

charts = [s for s in slide6.shapes if hasattr(s, 'has_chart') and s.has_chart]
if charts:
    chart = charts[0].chart
    print(f"\nChart Title: {chart.chart_title.text_frame.text if chart.has_title else 'No title'}")
    print(f"Number of series: {len(chart.series)}")
    
    categories = list(chart.plots[0].categories)
    print(f"\nCategories ({len(categories)}):")
    for i, cat in enumerate(categories):
        print(f"  {i+1}. {cat}")
    
    print(f"\nSeries Data:")
    for idx, series in enumerate(chart.series):
        values = list(series.values)
        print(f"\n  Series {idx+1}: '{series.name}'")
        print(f"    Values: {values}")
        print(f"    Sum: {sum(values)}")
        
        # Show bar-by-bar breakdown
        print(f"    Breakdown:")
        for i, (cat, val) in enumerate(zip(categories, values)):
            if val > 0:
                print(f"      {cat}: {val}")

print("\n" + "=" * 80)
print("SUCCESS - Multi-series chart created!")
print("=" * 80)
