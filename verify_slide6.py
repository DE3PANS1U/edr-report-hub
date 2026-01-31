"""
Quick verification of Slide 6 chart structure
"""
from pptx import Presentation

generated = Presentation("EDR-Weekly-Incident-Report_2026-01-31.pptx")
template = Presentation("EDR Weekly Incident Metrics (01st Dec 25 to 07th Dec 25).pptx")

print("=" * 80)
print("SLIDE 6 CHART COMPARISON")
print("=" * 80)

# Template
template_slide6 = template.slides[5]
template_charts = [s for s in template_slide6.shapes if hasattr(s, 'has_chart') and s.has_chart]
template_chart = template_charts[0].chart
template_values = list(template_chart.series[0].values)
template_categories = list(template_chart.plots[0].categories)

print("\nTEMPLATE:")
print(f"  Total bars: {len(template_values)}")
print(f"  Values: {template_values}")
print(f"  Categories: {template_categories}")

# Generated
generated_slide6 = generated.slides[5]
generated_charts = [s for s in generated_slide6.shapes if hasattr(s, 'has_chart') and s.has_chart]
generated_chart = generated_charts[0].chart
generated_values = list(generated_chart.series[0].values)
generated_categories = list(generated_chart.plots[0].categories)

print("\nGENERATED:")
print(f"  Total bars: {len(generated_values)}")
print(f"  Values: {generated_values}")
print(f"  Categories: {generated_categories}")

print("\n" + "=" * 80)
print("COMPARISON:")
print(f"  Template bars: {len(template_values)}")
print(f"  Generated bars: {len(generated_values)}")
print(f"  Match: {'✓ YES' if len(template_values) == len(generated_values) else '✗ NO'}")
print("=" * 80)
