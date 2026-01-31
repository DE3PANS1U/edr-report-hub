"""
Check template chart data for slides 4-8
"""
from pptx import Presentation

template = "EDR Weekly Incident Metrics (01st Dec 25 to 07th Dec 25).pptx"
prs = Presentation(template)

print("TEMPLATE FILE CHART DATA ANALYSIS")
print("=" * 80)

for slide_idx in [3, 4, 5, 6, 7]:  # Slides 4-8 (0-indexed)
    slide = prs.slides[slide_idx]
    print(f"\n{'=' * 80}")
    print(f"TEMPLATE SLIDE {slide_idx + 1}")
    print(f"{'=' * 80}")
    
    chart_found = False
    for shape in slide.shapes:
        if hasattr(shape, "has_chart") and shape.has_chart:
            chart_found = True
            chart = shape.chart
            title = chart.chart_title.text_frame.text if chart.has_title else "No title"
            
            print(f"\nChart: {title}")
            print(f"  Type: {chart.chart_type}")
            
            try:
                for series in chart.series:
                    values = list(series.values)
                    print(f"  Series '{series.name}': {values}")
                    print(f"    Sum: {sum(values)}, Non-zero: {sum(1 for v in values if v > 0)}")
            except Exception as e:
                print(f"  ERROR: {e}")
    
    if not chart_found:
        print("  No charts found on this slide")

print("\n" + "=" * 80)
