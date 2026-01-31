"""
Comprehensive analysis of all slides in generated PowerPoint
"""
from pptx import Presentation
import os

pptx_file = "EDR-Weekly-Incident-Report_2026-01-31.pptx"

if not os.path.exists(pptx_file):
    print(f"ERROR: File not found: {pptx_file}")
    exit(1)

prs = Presentation(pptx_file)

print("=" * 80)
print(f"ANALYZING: {pptx_file}")
print(f"Total Slides: {len(prs.slides)}")
print(f"Dimensions: {prs.slide_width.inches:.2f}\" x {prs.slide_height.inches:.2f}\"")
print("=" * 80)

for slide_idx, slide in enumerate(prs.slides, 1):
    print(f"\n{'=' * 80}")
    print(f"SLIDE {slide_idx}")
    print(f"{'=' * 80}")
    
    # Count different shape types
    text_boxes = 0
    charts = 0
    tables = 0
    images = 0
    other_shapes = 0
    
    print(f"Total shapes on slide: {len(slide.shapes)}")
    
    for shape_idx, shape in enumerate(slide.shapes):
        shape_type = shape.shape_type
        print(f"\nShape {shape_idx + 1}:")
        print(f"  Type: {shape_type}")
        print(f"  Position: ({shape.left.inches:.2f}\", {shape.top.inches:.2f}\")")
        print(f"  Size: {shape.width.inches:.2f}\" x {shape.height.inches:.2f}\"")
        
        # Check for text
        if hasattr(shape, "text"):
            text = shape.text.strip()
            if text:
                text_boxes += 1
                print(f"  Text: {text[:100]}")
        
        # Check for chart
        if hasattr(shape, "has_chart") and shape.has_chart:
            charts += 1
            chart = shape.chart
            print(f"  Chart Type: {chart.chart_type}")
            if chart.has_title:
                print(f"  Chart Title: {chart.chart_title.text_frame.text}")
        
        # Check for table
        if hasattr(shape, "has_table") and shape.has_table:
            tables += 1
            table = shape.table
            print(f"  Table: {len(table.rows)} rows x {len(table.columns)} columns")
            # Print first few cells
            if len(table.rows) > 0 and len(table.columns) > 0:
                print(f"  First cell: {table.cell(0, 0).text}")
    
    print(f"\nSUMMARY FOR SLIDE {slide_idx}:")
    print(f"  Text boxes: {text_boxes}")
    print(f"  Charts: {charts}")
    print(f"  Tables: {tables}")
    print(f"  Total shapes: {len(slide.shapes)}")
    
    if len(slide.shapes) == 0:
        print("  *** WARNING: SLIDE IS COMPLETELY EMPTY! ***")
    elif text_boxes == 0 and charts == 0 and tables == 0:
        print("  *** WARNING: NO CONTENT DETECTED! ***")

print("\n" + "=" * 80)
print("ANALYSIS COMPLETE")
print("=" * 80)
