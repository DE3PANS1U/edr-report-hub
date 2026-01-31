"""Check what color 'closed' is currently using"""
from pptx import Presentation

generated = Presentation("EDR-Weekly-Incident-Report_2026-01-31.pptx")
template = Presentation("EDR Weekly Incident Metrics (01st Dec 25 to 07th Dec 25).pptx")

print("=" * 60)
print("SLIDE 3 - Alert Status Chart Colors")
print("=" * 60)

# Check generated
gen_slide3 = generated.slides[2]
gen_charts = [s for s in gen_slide3.shapes if hasattr(s, 'has_chart') and s.has_chart]
if gen_charts:
    chart = gen_charts[0].chart
    print("\nGENERATED - First chart:")
    try:
        for idx, point in enumerate(chart.series[0].points):
            cat = list(chart.plots[0].categories)[idx]
            rgb = point.format.fill.fore_color.rgb
            print(f"  {cat}: RGB({rgb[0]}, {rgb[1]}, {rgb[2]})")
    except:
        print("  Could not extract colors")

# Check template
temp_slide3 = template.slides[2]
temp_charts = [s for s in temp_slide3.shapes if hasattr(s, 'has_chart') and s.has_chart]
if temp_charts:
    chart = temp_charts[0].chart
    print("\nTEMPLATE - First chart:")
    try:
        for idx, point in enumerate(chart.series[0].points):
            cat = list(chart.plots[0].categories)[idx]
            rgb = point.format.fill.fore_color.rgb
            print(f"  {cat}: RGB({rgb[0]}, {rgb[1]}, {rgb[2]})")
    except Exception as e:
        print(f"  Error: {e}")

print("\n" + "=" * 60)
