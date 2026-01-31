"""
Analyze existing PowerPoint template and Excel data
"""
import pandas as pd
from pptx import Presentation
import os

# Read combined.xlsx
print("=" * 60)
print("ANALYZING COMBINED.XLSX")
print("=" * 60)

df = pd.read_excel('combined.xlsx')
print(f"\nShape: {df.shape}")
print(f"\nColumn Names:")
for i, col in enumerate(df.columns):
    print(f"  {i+1}. {col}")

print(f"\nFirst 15 rows:")
print(df.head(15).to_string())

print(f"\nData Types:")
print(df.dtypes)

print(f"\nUnique Values Per Column:")
for col in df.columns:
    unique_vals = df[col].unique()
    print(f"\n{col}: {len(unique_vals)} unique values")
    if len(unique_vals) <= 20:
        print(f"  Values: {unique_vals.tolist()}")

# Read existing PowerPoint
print("\n" + "=" * 60)
print("ANALYZING EXISTING POWERPOINT TEMPLATE")
print("=" * 60)

ppt_file = "EDR Weekly Incident Metrics (01st Dec 25 to 07th Dec 25).pptx"
if os.path.exists(ppt_file):
    prs = Presentation(ppt_file)
    
    print(f"\nTotal Slides: {len(prs.slides)}")
    print(f"Slide Dimensions: {prs.slide_width.inches} x {prs.slide_height.inches} inches")
    
    for slide_idx, slide in enumerate(prs.slides):
        print(f"\n{'=' * 50}")
        print(f"SLIDE {slide_idx + 1}")
        print(f"{'=' * 50}")
        print(f"Layout: {slide.slide_layout.name}")
        
        # Extract text from shapes
        for shape_idx, shape in enumerate(slide.shapes):
            if hasattr(shape, "text") and shape.text.strip():
                print(f"\nShape {shape_idx + 1} - {shape.shape_type}: {shape.text[:100]}")
            
            # Check if it's a chart
            if hasattr(shape, "has_chart") and shape.has_chart:
                chart = shape.chart
                print(f"\nChart Found - Type: {chart.chart_type}")
                print(f"  Title: {chart.chart_title.text_frame.text if chart.has_title else 'No title'}")
                
            # Check if it's a table
            if hasattr(shape, "has_table") and shape.has_table:
                table = shape.table
                print(f"\nTable Found - {len(table.rows)} rows x {len(table.columns)} columns")
                # Print first few rows
                for r_idx in range(min(3, len(table.rows))):
                    row_data = []
                    for c_idx in range(len(table.columns)):
                        cell_text = table.cell(r_idx, c_idx).text
                        row_data.append(cell_text[:30])
                    print(f"  Row {r_idx}: {row_data}")
else:
    print("\nPowerPoint file not found!")
